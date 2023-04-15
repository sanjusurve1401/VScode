
import os
os.environ["IMAGEIO_FFMPEG_EXE"] = "/path/to/ffmpeg"

import tkinter
from tkinter import filedialog
import zipfile
import PyPDF2
from pptx import Presentation
import moviepy.editor as mp

# Define the dimensions of the video frames
frame_width = 640
frame_height = 480

# Create a list to hold the file paths of converted image files
image_paths = []

# Create a temporary folder to store the converted image files
temp_folder = "temp"
os.makedirs(temp_folder, exist_ok=True)

# Create a Tkinter root window to allow the user to select a zip file
root = tkinter.Tk()
root.withdraw()

# Ask the user to select a zip file and get the file path
file_path = filedialog.askopenfilename()

# Extract the contents of the zip file into the temporary folder
zip_file = zipfile.ZipFile(file_path)
zip_file.extractall(temp_folder)

# Loop through each file in the temporary folder and convert PDFs and PPTs to images
for file_name in os.listdir(temp_folder):
    file_path = os.path.join(temp_folder, file_name)

    # Check if the file is a PDF file
    if file_name.lower().endswith(".pdf"):
        # Open the PDF file using PyPDF2
        with open(file_path, "rb") as pdf_file:
            pdf_reader = PyPDF2.PdfReader(pdf_file)

            # Loop through each page in the PDF file and convert it to an image
            for page_num in range(pdf_reader.getNumPages()):
                pdf_page = pdf_reader.getPage(page_num)
                pdf_page_img = pdf_page.convertTo("png")
                pdf_page_img_path = os.path.join(temp_folder, f"{file_name}_{page_num}.png")
                pdf_page_img.save(pdf_page_img_path)
                image_paths.append(pdf_page_img_path)

    # Check if the file is a PPT file
    elif file_name.lower().endswith((".ppt", ".pptx")):
        # Open the PPT file using python-pptx
        ppt = Presentation(file_path)

        # Loop through each slide in the PPT file and convert it to an image
        for i, slide in enumerate(ppt.slides):
            slide_img_path = os.path.join(temp_folder, f"{file_name}_{i}.png")
            slide.shapes.export(slide_img_path)
            image_paths.append(slide_img_path)

# Sort the image paths based on their filename
image_paths.sort()

# Create a list of moviepy ImageClips from the image paths
clips = [mp.ImageClip(img_path).set_duration(2)
         for img_path in image_paths]

# Concatenate the ImageClips into a single video clip
video_clip = mp.concatenate_videoclips(clips)

# Set the video clip dimensions and write it to a file
video_clip = video_clip.resize((frame_width, frame_height))
video_clip.write_videofile("output.mp4", fps=24)